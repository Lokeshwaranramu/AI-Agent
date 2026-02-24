"""
Document modification tools for Word, Excel, and PowerPoint files.
"""
import os
from pathlib import Path
from typing import Any, Dict, List, Optional

import openpyxl
from docx import Document
from pptx import Presentation
from pptx.util import Pt as PptPt

from utils.logger import log


def modify_word_document(
    input_path: str,
    replacements: Dict[str, str],
    output_path: Optional[str] = None,
    append_text: Optional[str] = None,
) -> str:
    """
    Modify a Word document: find/replace text, append content.

    Args:
        input_path: Path to the .docx file
        replacements: Dict of {old_text: new_text} replacements
        output_path: Output file path (modifies in place if None)
        append_text: Optional text to append at the end

    Returns:
        Path to the modified document
    """
    if not Path(input_path).exists():
        raise FileNotFoundError(f"Document not found: {input_path}")

    if output_path is None:
        output_path = input_path

    try:
        doc = Document(input_path)

        for paragraph in doc.paragraphs:
            for old_text, new_text in replacements.items():
                if old_text in paragraph.text:
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            log.debug(f"Replaced '{old_text}' with '{new_text}'")

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for old_text, new_text in replacements.items():
                        if old_text in cell.text:
                            cell.text = cell.text.replace(old_text, new_text)

        if append_text:
            doc.add_paragraph(append_text)

        os.makedirs(Path(output_path).parent, exist_ok=True)
        doc.save(output_path)
        log.success(f"Document modified and saved: {output_path}")
        return output_path

    except Exception as e:
        log.error(f"Document modification failed: {e}")
        raise


def create_word_document(content: Dict[str, Any], output_path: str) -> str:
    """
    Create a new Word document from structured content.

    Args:
        content: Dict with keys: title, sections (list of {heading, body, bullet_points})
        output_path: Where to save the document

    Returns:
        Path to the created document
    """
    try:
        doc = Document()

        if "title" in content:
            doc.add_heading(content["title"], level=0)

        for section in content.get("sections", []):
            if section.get("heading"):
                doc.add_heading(section["heading"], level=1)
            if section.get("body"):
                doc.add_paragraph(section["body"])
            for point in section.get("bullet_points", []):
                doc.add_paragraph(point, style="List Bullet")

        os.makedirs(Path(output_path).parent, exist_ok=True)
        doc.save(output_path)
        log.success(f"Word document created: {output_path}")
        return output_path

    except Exception as e:
        log.error(f"Document creation failed: {e}")
        raise


def modify_excel_file(
    input_path: str,
    cell_updates: Dict[str, Any],
    sheet_name: Optional[str] = None,
    output_path: Optional[str] = None,
) -> str:
    """
    Modify an Excel file: update specific cells.

    Args:
        input_path: Path to .xlsx file
        cell_updates: Dict of {cell_address: value} e.g. {'A1': 'Hello', 'B2': 42}
        sheet_name: Sheet to modify (first sheet if None)
        output_path: Output path (modifies in place if None)

    Returns:
        Path to the modified Excel file
    """
    if not Path(input_path).exists():
        raise FileNotFoundError(f"Excel file not found: {input_path}")

    if output_path is None:
        output_path = input_path

    try:
        wb = openpyxl.load_workbook(input_path)
        ws = wb[sheet_name] if sheet_name else wb.active

        for cell_addr, value in cell_updates.items():
            ws[cell_addr] = value
            log.debug(f"Updated cell {cell_addr} = {value}")

        os.makedirs(Path(output_path).parent, exist_ok=True)
        wb.save(output_path)
        log.success(f"Excel file modified: {output_path}")
        return output_path

    except Exception as e:
        log.error(f"Excel modification failed: {e}")
        raise


def modify_powerpoint(
    input_path: str,
    slide_text_replacements: Dict[str, str],
    output_path: Optional[str] = None,
) -> str:
    """
    Modify a PowerPoint file: find/replace text across all slides.

    Args:
        input_path: Path to .pptx file
        slide_text_replacements: Dict of {old_text: new_text}
        output_path: Output path (modifies in place if None)

    Returns:
        Path to the modified presentation
    """
    if not Path(input_path).exists():
        raise FileNotFoundError(f"PowerPoint not found: {input_path}")

    if output_path is None:
        output_path = input_path

    try:
        prs = Presentation(input_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for old_text, new_text in slide_text_replacements.items():
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)

        os.makedirs(Path(output_path).parent, exist_ok=True)
        prs.save(output_path)
        log.success(f"PowerPoint modified: {output_path}")
        return output_path

    except Exception as e:
        log.error(f"PowerPoint modification failed: {e}")
        raise
